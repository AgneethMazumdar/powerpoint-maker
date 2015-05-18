# Released under MIT License 
# Created By Agneeth Mazumdar

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import urllib
import os
import csv

def read_csv():

    names = []
    urls = []
    
    with open('your_csv_file_here.csv', 'rb') as names_images_data:
        csv_names_images_data = csv.reader(names_images_data)

        for row in csv_names_images_data:
            names.append(row[0])
            urls.append(row[1])

        names_plus_images = dict(zip(names, urls))

    return names, urls, names_plus_images

def add_jpg_to_image_names(names_list):
	
    extended_names = names_list 

    for index, val in enumerate(extended_names):
        if ' ' in extended_names[index]:
            extended_names[index] = extended_names[index].replace(' ', '');
            extended_names[index] = extended_names[index] + '.jpg'
        else:
            extended_names[index] = extended_names[index] + '.jpg'
    
    return extended_names
        
def download_images(names_list, urls_list):

    names = names_list 
    urls = urls_list

    for index, val in enumerate(names):
        urllib.urlretrieve(urls[index], names[index]) 

        im = Image.open(names[index]) 
        out = im.thumbnail((400, 380), Image.ANTIALIAS)  
        im.save(names[index], "JPEG") 

    return names

def original_names(dictionary):
    
    names_plus_images = dictionary

    preserved_names = []

    for index, key in enumerate(names_plus_images):
        preserved_names.append(key)

    return preserved_names 

def make_slides(powerpoint, names_list, original_names, counter):

    prs = powerpoint

    names = names_list 
    old_names = original_names
    index = counter
       
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = old_names[index]
    image_name = names[index]

    left = Inches(2.5)
    top = Inches(1.7)

    pic = slide.shapes.add_picture(image_name, left, top)

    prs.save('insert_desired_name_here.pptx')

def main():

    prs = Presentation()

    names, urls, names_plus_images = read_csv()
    
    new_names = add_jpg_to_image_names(names)
    download_images(new_names, urls)
    old_names = original_names(names_plus_images)

    retrieved_names, urls, names_plus_images = read_csv()

    for index, val in enumerate(names):
        make_slides(prs, names, retrieved_names, index)

if __name__ == "__main__":
    main()
